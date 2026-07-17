#!/usr/bin/env python3
"""
render.py — fill a branded PowerPoint template from a Markdown artifact.

This is the deterministic renderer for the `markdown-to-pptx` skill. The
agent assembles nothing here: it invokes this script, which inspects a
user-provided `.pptx` template's layout placeholders and projects a Markdown
document onto them. The user's slide master, theme, and placed assets survive
because we fill an existing template rather than building a deck from scratch.

Verbs
    --check                 import-probe python-pptx; exit 0 (present) / 2 (absent)
    inspect <template>      print the layout/placeholder manifest
    render  <markdown>      fill a template (or the library default) and write a .pptx
        --template <path>   branded template to fill; omit for the opt-out default
        --output <path>     output path (default: <markdown-basename>.pptx in CWD)

Stdout markers (for agent parsing)
    FILLPOINTS: layout=<i> idx=<i> type=<NAME> name=<text>   (inspect)
    OUTPUT: <path>          path to the written .pptx
    FILLED: <n>             count of placeholders filled
    WARNING: <msg>          non-fatal issue, surface to the user

Trust model
    A user-supplied template is treated as trusted-author input, consistent
    with the converters pack's local-files-trusted stance. python-pptx does not
    evaluate template content as code, so there is no SSTI surface here; XXE /
    zip-bomb on a crafted Office archive is an accepted, out-of-scope risk for a
    trusted-author template. Output-path confinement (below) is still enforced
    because the output path is assembled from model-influenced content.

Errors go to stderr; exit code 1 on failure, 2 on a missing library.
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

PIP_INSTALL = "python -m pip install 'python-pptx>=1.0.0'"


# --------------------------------------------------------------------------- #
# Output-path confinement
# --------------------------------------------------------------------------- #
def confine(path: Path, root: Path) -> Path:
    """Resolve ``path`` (following symlinks) and require it under ``root``.

    Uses a path-*component* containment check (``root`` is the resolved path or
    appears among its ``.parents``), not a string-prefix check, so a sibling
    like ``workdir-evil`` is rejected against root ``workdir``. Raises
    ``ValueError`` on escape.
    """
    resolved = path.resolve()
    root_resolved = root.resolve()
    if resolved == root_resolved or root_resolved in resolved.parents:
        return resolved
    raise ValueError(
        f"path {path!r} resolves to {resolved} which is outside the working "
        f"directory {root_resolved}; refusing to read/write outside the project"
    )


# --------------------------------------------------------------------------- #
# Markdown → content model
# --------------------------------------------------------------------------- #
def parse_markdown(text: str) -> dict:
    """Project a Markdown artifact onto the pptx content model.

    Mapping: front-matter ``key: value`` lines → ``scalars``; each
    H1/H2 heading → one ``section`` (→ one slide); list items under a heading →
    ``bullets``; a Markdown table under a heading → ``table`` (rows of cells).

    Returns ``{"scalars": {...}, "sections": [{"heading", "bullets", "table"}]}``.
    Pure transform — no I/O.
    """
    scalars: dict[str, str] = {}
    sections: list[dict] = []
    lines = text.splitlines()
    i = 0

    # YAML-ish front-matter: a leading `---` fence of simple `key: value` lines.
    if lines and lines[0].strip() == "---":
        j = 1
        while j < len(lines) and lines[j].strip() != "---":
            raw = lines[j]
            if ":" in raw:
                key, _, value = raw.partition(":")
                scalars[key.strip()] = value.strip().strip("\"'")
            j += 1
        i = j + 1 if j < len(lines) else len(lines)

    current: dict | None = None

    def cell_split(row: str) -> list[str]:
        parts = row.strip().strip("|").split("|")
        return [c.strip() for c in parts]

    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        if stripped.startswith("#"):
            level = len(stripped) - len(stripped.lstrip("#"))
            if level in (1, 2):
                current = {"heading": stripped[level:].strip(), "bullets": [], "table": None}
                sections.append(current)
                i += 1
                continue

        if current is not None and (stripped.startswith("- ") or stripped.startswith("* ")):
            current["bullets"].append(stripped[2:].strip())
            i += 1
            continue

        # Markdown table: a `|`-delimited row whose successor is a `---` divider.
        if (
            current is not None
            and stripped.startswith("|")
            and i + 1 < len(lines)
            and "-" in lines[i + 1]
            and set(lines[i + 1].strip().replace("|", "").replace(":", "").strip()) <= {"-", " "}
            and lines[i + 1].strip().replace("|", "").strip()
        ):
            header = cell_split(stripped)
            rows = [header]
            i += 2  # skip header + divider
            while i < len(lines) and lines[i].strip().startswith("|"):
                rows.append(cell_split(lines[i]))
                i += 1
            current["table"] = rows
            continue

        # Plain prose line under a heading → a body paragraph (bullet without marker).
        if current is not None and stripped:
            current["bullets"].append(stripped)
        i += 1

    return {"scalars": scalars, "sections": sections}


# --------------------------------------------------------------------------- #
# python-pptx helpers (imported lazily so --check can run without the library)
# --------------------------------------------------------------------------- #
def _placeholder_type_name(ph) -> str:
    try:
        return str(ph.placeholder_format.type).split()[0].split(".")[-1]
    except Exception:  # pragma: no cover — defensive
        return "UNKNOWN"


def inspect_template(template: Path) -> list[dict]:
    """Return the layout/placeholder manifest for a .pptx template.

    One record per placeholder across every slide layout, keyed by the layout
    index and the placeholder ``idx`` (its stable identity — not list position).
    """
    from pptx import Presentation

    prs = Presentation(str(template))
    manifest: list[dict] = []
    for layout_i, layout in enumerate(prs.slide_layouts):
        for ph in layout.placeholders:
            manifest.append(
                {
                    "layout": layout_i,
                    "idx": ph.placeholder_format.idx,
                    "type": _placeholder_type_name(ph),
                    "name": ph.name,
                }
            )
    return manifest


def _find_layout(prs, *type_names):
    """First slide layout that carries a placeholder of any of ``type_names``."""
    for layout in prs.slide_layouts:
        for ph in layout.placeholders:
            if _placeholder_type_name(ph) in type_names:
                return layout
    return prs.slide_layouts[0]


def _placeholder_by_type(slide, *type_names):
    for ph in slide.placeholders:
        if _placeholder_type_name(ph) in type_names:
            return ph
    return None


def render_pptx(model: dict, template: Path | None, output: Path) -> tuple[int, list[str]]:
    """Fill ``template`` (or the library default) from ``model``; write ``output``.

    Returns ``(filled_count, warnings)``.
    """
    from pptx import Presentation
    from pptx.util import Inches

    warnings: list[str] = []
    prs = Presentation(str(template)) if template is not None else Presentation()
    if template is None:
        warnings.append(
            "no template supplied — rendering with the python-pptx default master; "
            "the result carries no brand"
        )
    filled = 0

    # Title slide from front-matter scalars.
    scalars = model.get("scalars", {})
    if scalars.get("title") or scalars.get("subtitle"):
        title_layout = _find_layout(prs, "TITLE", "CENTER_TITLE")
        slide = prs.slides.add_slide(title_layout)
        title_ph = _placeholder_by_type(slide, "TITLE", "CENTER_TITLE")
        if title_ph is not None and scalars.get("title"):
            title_ph.text = scalars["title"]
            filled += 1
        sub_ph = _placeholder_by_type(slide, "SUBTITLE", "BODY")
        if sub_ph is not None and scalars.get("subtitle"):
            sub_ph.text = scalars["subtitle"]
            filled += 1

    content_layout = _find_layout(prs, "BODY", "OBJECT")
    for section in model.get("sections", []):
        slide = prs.slides.add_slide(content_layout)
        title_ph = _placeholder_by_type(slide, "TITLE", "CENTER_TITLE")
        if title_ph is not None:
            title_ph.text = section.get("heading", "")
            filled += 1

        bullets = section.get("bullets") or []
        if bullets:
            body = _placeholder_by_type(slide, "BODY", "OBJECT")
            if body is not None and body.has_text_frame:
                tf = body.text_frame
                tf.text = bullets[0]
                for b in bullets[1:]:
                    tf.add_paragraph().text = b
                filled += 1
            else:
                warnings.append(
                    f"slide '{section.get('heading')}' has bullets but the layout "
                    "exposes no body placeholder; bullets dropped"
                )

        table = section.get("table")
        if table:
            # Fill a TABLE placeholder when the layout has one; otherwise add a
            # table graphic frame to the slide. python-pptx returns a *new*
            # graphic-frame object from insert_table — re-fetch, never reuse.
            table_ph = _placeholder_by_type(slide, "TABLE")
            rows, cols = len(table), max(len(r) for r in table)
            if table_ph is not None:
                gframe = table_ph.insert_table(rows=rows, cols=cols)
                tbl = gframe.table
            else:
                gframe = slide.shapes.add_table(
                    rows, cols, Inches(0.5), Inches(2.0), Inches(9.0), Inches(0.4 * rows)
                )
                tbl = gframe.table
            for r, row in enumerate(table):
                for c in range(cols):
                    tbl.cell(r, c).text = row[c] if c < len(row) else ""
            filled += 1

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    return filled, warnings


# --------------------------------------------------------------------------- #
# CLI
# --------------------------------------------------------------------------- #
def cmd_check() -> int:
    try:
        import pptx  # noqa: F401
    except ImportError:
        print(f"python-pptx is not installed. Run:\n    {PIP_INSTALL}", file=sys.stderr)
        return 2
    return 0


def cmd_inspect(args) -> int:
    root = Path.cwd()
    template = confine(Path(args.template), root)
    if not template.is_file():
        print(f"ERROR: template not found: {args.template}", file=sys.stderr)
        return 1
    for rec in inspect_template(template):
        print(
            f"FILLPOINTS: layout={rec['layout']} idx={rec['idx']} "
            f"type={rec['type']} name={rec['name']}"
        )
    return 0


def cmd_render(args) -> int:
    root = Path.cwd()
    md_path = confine(Path(args.markdown), root)
    if not md_path.is_file():
        print(f"ERROR: markdown not found: {args.markdown}", file=sys.stderr)
        return 1

    template = confine(Path(args.template), root) if args.template else None
    if template is not None and not template.is_file():
        print(f"ERROR: template not found: {args.template}", file=sys.stderr)
        return 1

    output = Path(args.output) if args.output else Path(md_path.stem + ".pptx")
    output = confine(output, root)

    model = parse_markdown(md_path.read_text(encoding="utf-8"))
    filled, warnings = render_pptx(model, template, output)

    for w in warnings:
        print(f"WARNING: {w}")
    print(f"OUTPUT: {output}")
    print(f"FILLED: {filled}")
    return 0


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Fill a PowerPoint template from Markdown.")
    parser.add_argument("--check", action="store_true", help="probe python-pptx; exit 0/2")
    sub = parser.add_subparsers(dest="verb")

    p_inspect = sub.add_parser("inspect", help="print the placeholder manifest")
    p_inspect.add_argument("template")

    p_render = sub.add_parser("render", help="fill a template and write a .pptx")
    p_render.add_argument("markdown")
    p_render.add_argument("--template", help="branded .pptx to fill (omit for the default)")
    p_render.add_argument("--output", help="output path (default: <basename>.pptx in CWD)")
    return parser


def main(argv: list[str] | None = None) -> int:
    parser = build_parser()
    args = parser.parse_args(argv)

    if args.check:
        return cmd_check()

    # Every other verb needs the library present.
    rc = cmd_check()
    if rc != 0:
        return rc

    try:
        if args.verb == "inspect":
            return cmd_inspect(args)
        if args.verb == "render":
            return cmd_render(args)
    except ValueError as exc:  # confinement refusal
        print(f"ERROR: {exc}", file=sys.stderr)
        return 1

    parser.print_help()
    return 1


if __name__ == "__main__":
    raise SystemExit(main())
