#!/usr/bin/env python3
"""
convert.py — convert documents and images to Markdown, tiered and no-ML-first.

The skill's document surface. It routes each input to the lowest tier that can
handle it:

  * **Tier 0 (no ML)** — pure-Python / stdlib extractors for digital PDFs
    (``pypdf``), Office files (``python-docx`` / ``openpyxl`` / ``python-pptx``,
    degrading to stdlib ``zipfile`` + XML), and the everyday text formats
    (HTML, EPUB, CSV/TSV, OpenDocument, ``.eml``). This is the floor a
    locked-down environment can always reach.
  * **Tier 2 (approved ML)** — Docling, unchanged, as the higher-fidelity
    fall-through for ``.xls`` and images. Its Markdown body is passed through
    to the contract builder unmodified.

Every path emits the one versioned unified output contract (frontmatter) via
``contract.build_frontmatter`` and writes through an output-path confinement
guard. Untrusted input is parsed defensively (see ``safe_io``): XXE-safe XML,
decompression-bomb guards, and a coarse resource ceiling.

Usage:
    python scripts/convert.py <file> [file2 ...]
    python scripts/convert.py --check [library ...]   # probe optional Tier-0 libs

Output:
    <input-basename>.md written next to the input file, confined to its dir.

Stdout markers (for agent parsing):
    OUTPUT: <path>   — path to the written Markdown file
    LINES: <n>       — line count of the output
    WORDS: <n>       — word count of the output
    WARNING: <msg>   — non-fatal issue (e.g. requires-review / escalation)

Errors go to stderr; exit code 1 on failure, 2 on a missing probed library.
"""
from __future__ import annotations

import argparse
import csv as csvmod
import io
import sys
import tempfile
import zipfile
from pathlib import Path
from typing import Callable, NamedTuple

import contract
import safe_io

# --- Format routing ---------------------------------------------------------

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp", ".webp", ".gif"}
# Handled only by Docling (Tier 2): legacy .xls (binary, no clean stdlib path)
# and images (OCR / vision).
DOCLING_EXTS = {".xls"} | IMAGE_EXTS

# Images wider or taller than this are pre-scaled before Docling processes them.
MAX_IMAGE_DIM = 4000

# Sparse-text threshold: a digital PDF yielding fewer than this many words
# is treated as image-only / low-quality and escalated, not silently emitted.
SPARSE_WORD_THRESHOLD = 20

# Coarse per-parser ceilings.
MAX_PDF_PAGES = 5_000
MAX_SHEET_ROWS = 1_000_000
MAX_CSV_ROWS = 1_000_000

# Optional Tier-0 libraries, resolved pip-on-demand via --check (never
# auto-installed), mirroring the sibling markdown-to-* skills.
OPTIONAL_LIBS = {
    "pypdf": "python -m pip install 'pypdf>=4.0.0'",
    "docx": "python -m pip install 'python-docx>=1.1.0'",       # import name: docx
    "openpyxl": "python -m pip install 'openpyxl>=3.1.0'",
    "pptx": "python -m pip install 'python-pptx>=1.0.0'",       # import name: pptx
}


class ExtractResult(NamedTuple):
    body: str                 # Markdown body (goes below the frontmatter fence)
    tier: str                 # contract.TIER_*
    content_type: str
    confidence: str           # high | medium | low
    requires_review: bool
    escalation: str | None = None   # e.g. contract.TIER_1 when text is sparse


# --- Extractor registry -----------------------------------------------------
# Populated at the bottom of the module once every extractor is defined.
_EXTRACTORS: dict[str, Callable[[Path], ExtractResult]] = {}


def dispatch(path: Path, *, enrich: bool = False) -> ExtractResult:
    """Route an input to its Tier-0 extractor, or fall through to Docling.

    ``enrich`` only affects the Docling (Tier-2) fall-through; Tier-0 extractors
    ignore it. This function constructs only Tiers 0–2 — Tier 3 is never reachable
    from here (it is produced solely by ``tier3.assemble_tier3`` via ``--tier3``)."""
    extractor = _EXTRACTORS.get(path.suffix.lower())
    if extractor is not None:
        return extractor(path)
    return _extract_docling(path, enrich=enrich)


def supported_exts() -> set[str]:
    return set(_EXTRACTORS) | DOCLING_EXTS


# --- Shared guards ----------------------------------------------------------


def _refused_result(content_type: str, message: str) -> ExtractResult:
    """A refused-but-flagged result: the input was not fully parsed — a resource
    ceiling or a defensive-parse refusal (a DTD, a decompression bomb). The
    output carries requires-review + the reason rather than passing silently
    or crashing the batch."""
    return ExtractResult(
        body=f"> **Not extracted.** {message}\n",
        tier=contract.TIER_0,
        content_type=content_type,
        confidence="low",
        requires_review=True,
    )


# Two call-site vocabularies (resource-ceiling vs defensive-parse refusal), one
# flagged shape — aliased so the two intents read distinctly at the call site
# without duplicating the body.
_ceiling_result = _refused_result
_defensive_result = _refused_result


def _guard_size(path: Path, content_type: str) -> ExtractResult | None:
    try:
        safe_io.check_input_size(path)
    except safe_io.ResourceCeilingError as exc:
        return _ceiling_result(content_type, str(exc))
    return None


def _lib_available(import_name: str) -> bool:
    import importlib.util

    return importlib.util.find_spec(import_name) is not None


# --- Tier 0: PDF (pypdf) ----------------------------------------------------


def _assess_text(text: str) -> tuple[str, bool, str | None]:
    """Map extracted text to (confidence, requires_review, escalation).

    Sparse or empty text escalates honestly to Tier 1 (agent-vision) rather
    than emitting silent low-quality Markdown."""
    if len(text.split()) < SPARSE_WORD_THRESHOLD:
        return "low", True, contract.TIER_1
    return "high", False, None


def _extract_pdf(path: Path) -> ExtractResult:
    if (over := _guard_size(path, "pdf")) is not None:
        return over
    if not _lib_available("pypdf"):
        # No stdlib PDF text path — degrade to the sparse-text escalation path
        # rather than importing Docling or hard-failing.
        return ExtractResult(
            body=(
                "> **Tier-0 PDF text extraction needs `pypdf`, which is not "
                "installed.** Install it (`" + OPTIONAL_LIBS["pypdf"] + "`) for "
                "no-ML text extraction, or escalate to Tier 1 (agent-vision).\n"
            ),
            tier=contract.TIER_0,
            content_type="pdf",
            confidence="low",
            requires_review=True,
            escalation=contract.TIER_1,
        )

    from pypdf import PdfReader

    reader = PdfReader(str(path))
    n_pages = len(reader.pages)
    if n_pages > MAX_PDF_PAGES:
        return _ceiling_result(
            "pdf",
            f"PDF has {n_pages} pages, over the {MAX_PDF_PAGES}-page ceiling; "
            f"refusing to parse unbounded",
        )

    parts: list[str] = []
    for page in reader.pages:
        parts.append(page.extract_text() or "")
    text = "\n\n".join(p.strip() for p in parts if p.strip())

    confidence, requires_review, escalation = _assess_text(text)
    body = text if text.strip() else "> **No extractable text layer found.**\n"
    return ExtractResult(
        body=body,
        tier=contract.TIER_0,
        content_type="pdf",
        confidence=confidence,
        requires_review=requires_review,
        escalation=escalation,
    )


# --- Tier 0: Office (docx / xlsx / pptx) ------------------------------------
# Every Office file is opened through safe_io.open_safe_zip first (the
# decompression-bomb axes) and its XML members are DTD-gated, so no DTD ever
# reaches the ordinary library's transitive lxml parser — the stdlib-XXE-safe
# resolution applied to the library path.

_W = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
_S = "{http://schemas.openxmlformats.org/spreadsheetml/2006/main}"


def _open_office(path: Path, content_type: str):
    """Open an Office file with the decompression-bomb guards. Returns a SafeZip
    or an ExtractResult on a defensive refusal (bomb, or a file that is not a
    valid zip — e.g. a corrupt or mislabeled ``.docx``)."""
    if (over := _guard_size(path, content_type)) is not None:
        return over
    try:
        return safe_io.open_safe_zip(path)
    except safe_io.ZipBombError as exc:
        return _defensive_result(content_type, f"decompression-bomb guard: {exc}")
    except zipfile.BadZipFile as exc:
        return _defensive_result(content_type, f"not a valid Office (zip) file: {exc}")


def _harden_for_lib(sz: "safe_io.SafeZip", content_type: str) -> ExtractResult | None:
    """Before an ordinary Office library re-opens the raw file, fully validate
    every member through SafeZip — the per-member + cumulative decompression
    caps and the whole-buffer DTD refusal that the library path would otherwise
    bypass. Returns an ExtractResult on refusal, else None."""
    try:
        sz.harden_untrusted()
    except safe_io.ZipBombError as exc:
        return _defensive_result(content_type, f"decompression-bomb guard: {exc}")
    except safe_io.XmlSafetyError as exc:
        return _defensive_result(content_type, f"XML safety guard: {exc}")
    return None


def _local(tag: str) -> str:
    return tag.rsplit("}", 1)[-1]


def _extract_docx(path: Path) -> ExtractResult:
    opened = _open_office(path, "docx")
    if isinstance(opened, ExtractResult):
        return opened
    sz = opened
    try:
        if _lib_available("docx"):
            if (guard := _harden_for_lib(sz, "docx")) is not None:
                return guard
            import docx

            document = docx.Document(str(path))
            lines = [p.text for p in document.paragraphs if p.text.strip()]
            for table in document.tables:
                for row in table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        lines.append(" | ".join(cells))
            body = "\n\n".join(lines)
            confidence = "high"
        else:
            xml = sz.read_member("word/document.xml")
            root = safe_io.parse_xml(xml)
            paras: list[str] = []
            for p in root.iter(f"{_W}p"):
                texts = [t.text or "" for t in p.iter(f"{_W}t")]
                joined = "".join(texts).strip()
                if joined:
                    paras.append(joined)
            body = "\n\n".join(paras)
            confidence = "medium"
    finally:
        sz.close()
    # An empty extraction keeps requires_review=False deliberately: unlike a
    # sparse PDF (where an image-only page is the escalation signal), an empty
    # Office body emits a visible "No text found" placeholder a reader sees, and
    # a legitimately text-free doc should not be forced to review.
    body = body or "> **No text found in the document.**\n"
    return ExtractResult(body, contract.TIER_0, "docx", confidence, False)


def _extract_xlsx(path: Path) -> ExtractResult:
    opened = _open_office(path, "xlsx")
    if isinstance(opened, ExtractResult):
        return opened
    sz = opened
    truncated = False
    try:
        if _lib_available("openpyxl"):
            if (guard := _harden_for_lib(sz, "xlsx")) is not None:
                return guard
            import openpyxl

            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            sections: list[str] = []
            for ws in wb.worksheets:
                rows: list[str] = []
                for i, row in enumerate(ws.iter_rows(values_only=True)):
                    if i >= MAX_SHEET_ROWS:
                        rows.append(f"> _(truncated at {MAX_SHEET_ROWS} rows)_")
                        truncated = True
                        break
                    cells = ["" if c is None else str(c) for c in row]
                    if any(cells):
                        rows.append(" | ".join(cells))
                if rows:
                    sections.append(f"## {ws.title}\n\n" + "\n".join(rows))
            wb.close()
            body = "\n\n".join(sections)
            confidence = "high"
        else:
            body, confidence, truncated = _extract_xlsx_stdlib(sz)
    finally:
        sz.close()
    body = body or "> **No cell values found.**\n"
    # A row-truncated sheet is not fully extracted — flag it for review rather
    # than silently dropping the tail.
    if truncated:
        confidence = "low"
    return ExtractResult(body, contract.TIER_0, "xlsx", confidence, truncated)


def _extract_xlsx_stdlib(sz: "safe_io.SafeZip") -> tuple[str, str, bool]:
    # Shared strings table (cells with t="s" index into it).
    shared: list[str] = []
    if sz.has_member("xl/sharedStrings.xml"):
        root = safe_io.parse_xml(sz.read_member("xl/sharedStrings.xml"))
        for si in root.iter(f"{_S}si"):
            shared.append("".join(t.text or "" for t in si.iter(f"{_S}t")))
    sheet_names = sorted(
        n for n in sz.namelist()
        if n.startswith("xl/worksheets/") and n.endswith(".xml")
    )
    sections: list[str] = []
    truncated = False
    for name in sheet_names:
        root = safe_io.parse_xml(sz.read_member(name))
        rows: list[str] = []
        for i, row in enumerate(root.iter(f"{_S}row")):
            if i >= MAX_SHEET_ROWS:
                rows.append(f"> _(truncated at {MAX_SHEET_ROWS} rows)_")
                truncated = True
                break
            cells: list[str] = []
            for c in row.iter(f"{_S}c"):
                v = c.find(f"{_S}v")
                text = v.text if v is not None else None
                if text is None:
                    cells.append("")
                elif c.get("t") == "s":
                    idx = int(text)
                    cells.append(shared[idx] if 0 <= idx < len(shared) else "")
                else:
                    cells.append(text)
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            sections.append("\n".join(rows))
    return "\n\n".join(sections), "medium", truncated


def _extract_pptx(path: Path) -> ExtractResult:
    opened = _open_office(path, "pptx")
    if isinstance(opened, ExtractResult):
        return opened
    sz = opened
    try:
        if _lib_available("pptx"):
            if (guard := _harden_for_lib(sz, "pptx")) is not None:
                return guard
            import pptx

            prs = pptx.Presentation(str(path))
            slides: list[str] = []
            for i, slide in enumerate(prs.slides, 1):
                texts: list[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            line = "".join(run.text for run in para.runs).strip()
                            if line:
                                texts.append(line)
                slides.append(f"## Slide {i}\n\n" + "\n\n".join(texts))
            body = "\n\n".join(slides)
            confidence = "high"
        else:
            slide_names = sorted(
                n for n in sz.namelist()
                if n.startswith("ppt/slides/slide") and n.endswith(".xml")
            )
            slides = []
            for i, name in enumerate(slide_names, 1):
                root = safe_io.parse_xml(sz.read_member(name))
                texts = [t.text or "" for t in root.iter(f"{_A}t")]
                joined = "\n\n".join(t.strip() for t in texts if t.strip())
                slides.append(f"## Slide {i}\n\n" + joined)
            body = "\n\n".join(slides)
            confidence = "medium"
    finally:
        sz.close()
    body = body or "> **No text found in the presentation.**\n"
    return ExtractResult(body, contract.TIER_0, "pptx", confidence, False)


# --- Tier 0: D7 formats -----------------------------------------------------


class _TextHTMLParser:
    """Coarse HTML → text reduction via stdlib html.parser (no new dep).

    Emits block-level breaks for common structural tags and drops script/style
    content, which is enough for a context-layer floor (fidelity is Tier 2's
    job)."""

    def __init__(self) -> None:
        from html.parser import HTMLParser

        outer = self

        class _P(HTMLParser):
            def __init__(self) -> None:
                super().__init__(convert_charrefs=True)
                self.parts: list[str] = []
                self._skip = 0

            def handle_starttag(self, tag, attrs):
                if tag in ("script", "style"):
                    self._skip += 1
                elif tag in ("p", "div", "br", "li", "tr", "h1", "h2", "h3",
                             "h4", "h5", "h6", "section", "article"):
                    self.parts.append("\n")
                if tag == "li":
                    self.parts.append("- ")

            def handle_endtag(self, tag):
                if tag in ("script", "style") and self._skip:
                    self._skip -= 1
                elif tag in ("p", "div", "li", "tr", "h1", "h2", "h3", "h4",
                             "h5", "h6"):
                    self.parts.append("\n")

            def handle_data(self, data):
                if not self._skip and data.strip():
                    self.parts.append(data)

        self._parser = _P()

    def feed(self, text: str) -> str:
        self._parser.feed(text)
        raw = "".join(self._parser.parts)
        lines = [ln.strip() for ln in raw.splitlines()]
        out: list[str] = []
        for ln in lines:
            if ln or (out and out[-1]):
                out.append(ln)
        return "\n".join(out).strip()


def _extract_html(path: Path) -> ExtractResult:
    if (over := _guard_size(path, "html")) is not None:
        return over
    text = path.read_text(encoding="utf-8", errors="replace")
    body = _TextHTMLParser().feed(text) or "> **No text content found.**\n"
    return ExtractResult(body, contract.TIER_0, "html", "high", False)


def _extract_csv(path: Path) -> ExtractResult:
    content_type = "tsv" if path.suffix.lower() == ".tsv" else "csv"
    if (over := _guard_size(path, content_type)) is not None:
        return over
    delimiter = "\t" if content_type == "tsv" else ","
    text = path.read_text(encoding="utf-8", errors="replace")
    reader = csvmod.reader(io.StringIO(text), delimiter=delimiter)
    rows: list[list[str]] = []
    for i, row in enumerate(reader):
        if i >= MAX_CSV_ROWS:
            return _ceiling_result(
                content_type,
                f"{content_type.upper()} exceeds the {MAX_CSV_ROWS}-row ceiling; "
                f"refusing to parse unbounded",
            )
        rows.append(row)
    if not rows:
        return ExtractResult("> **Empty file.**\n", contract.TIER_0,
                             content_type, "high", False)
    ncol = max(len(r) for r in rows)
    header = rows[0] + [""] * (ncol - len(rows[0]))
    lines = ["| " + " | ".join(_md_cell(c) for c in header) + " |",
             "| " + " | ".join("---" for _ in header) + " |"]
    for r in rows[1:]:
        r = r + [""] * (ncol - len(r))
        lines.append("| " + " | ".join(_md_cell(c) for c in r) + " |")
    return ExtractResult("\n".join(lines), contract.TIER_0, content_type,
                         "high", False)


def _md_cell(s: str) -> str:
    return (s or "").replace("|", "\\|").replace("\n", " ").strip()


def _extract_epub(path: Path) -> ExtractResult:
    if (over := _guard_size(path, "epub")) is not None:
        return over
    try:
        sz = safe_io.open_safe_zip(path)
    except safe_io.ZipBombError as exc:
        return _defensive_result("epub", f"decompression-bomb guard: {exc}")
    try:
        # Read the spine order from the OPF; fall back to sorted xhtml members.
        xhtml = sorted(
            n for n in sz.namelist()
            if n.lower().endswith((".xhtml", ".html", ".htm"))
        )
        parser = _TextHTMLParser
        parts: list[str] = []
        dropped = False
        for name in xhtml:
            try:
                # DTD-gate each member, then reduce its (X)HTML to text.
                data = sz.read_member(name)
                safe_io._reject_dtd(data)
            except (safe_io.XmlSafetyError, safe_io.ZipBombError):
                dropped = True  # a guard skipped content — flag it, don't drop silently
                continue
            text = parser().feed(data.decode("utf-8", errors="replace"))
            if text:
                parts.append(text)
    finally:
        sz.close()
    body = "\n\n".join(parts) or "> **No readable text found in the EPUB.**\n"
    if dropped:
        body += ("\n\n> **Some content was skipped by a safety guard "
                 "(DTD or decompression-bomb) — review the source.**\n")
    return ExtractResult(body, contract.TIER_0, "epub",
                         "low" if dropped else "medium", dropped)


def _extract_odf(path: Path) -> ExtractResult:
    ct = {".odt": "odt", ".ods": "ods", ".odp": "odp"}[path.suffix.lower()]
    if (over := _guard_size(path, ct)) is not None:
        return over
    try:
        sz = safe_io.open_safe_zip(path)
    except safe_io.ZipBombError as exc:
        return _defensive_result(ct, f"decompression-bomb guard: {exc}")
    try:
        if not sz.has_member("content.xml"):
            return ExtractResult("> **No content.xml in the document.**\n",
                                 contract.TIER_0, ct, "low", True)
        root = safe_io.parse_xml(sz.read_member("content.xml"))
    except safe_io.XmlSafetyError as exc:
        return _defensive_result(ct, f"XML safety guard: {exc}")
    except safe_io.ZipBombError as exc:
        return _defensive_result(ct, f"decompression-bomb guard: {exc}")
    finally:
        sz.close()
    # OpenDocument text lives in <text:p>/<text:h>; join their text content.
    text_ns = "{urn:oasis:names:tc:opendocument:xmlns:text:1.0}"
    blocks: list[str] = []
    for tag in ("p", "h"):
        for el in root.iter(f"{text_ns}{tag}"):
            joined = "".join(el.itertext()).strip()
            if joined:
                blocks.append(joined)
    body = "\n\n".join(blocks) or "> **No text found in the document.**\n"
    return ExtractResult(body, contract.TIER_0, ct, "medium", False)


def _extract_eml(path: Path) -> ExtractResult:
    if (over := _guard_size(path, "eml")) is not None:
        return over
    import email
    from email import policy

    msg = email.message_from_bytes(path.read_bytes(), policy=policy.default)
    headers = []
    for h in ("From", "To", "Cc", "Date", "Subject"):
        if msg[h]:
            headers.append(f"**{h}:** {msg[h]}")
    try:
        part = msg.get_body(preferencelist=("plain", "html"))
    except Exception:
        part = None
    if part is not None:
        content = part.get_content()
        if part.get_content_type() == "text/html":
            content = _TextHTMLParser().feed(content)
    else:
        content = ""
    body = "\n\n".join(headers)
    if content and content.strip():
        body += "\n\n---\n\n" + content.strip()
    body = body or "> **No readable content in the message.**\n"
    return ExtractResult(body, contract.TIER_0, "eml", "high", False)


# --- Tier 2: Docling fall-through (unchanged body) --------------------------


def _prescale_image(input_path: Path, tmp_dir: str) -> Path:
    from PIL import Image

    Image.MAX_IMAGE_PIXELS = None  # disable PIL's bomb check for the resize step
    with Image.open(input_path) as img:
        w, h = img.size
        if max(w, h) <= MAX_IMAGE_DIM:
            return input_path
        scale = MAX_IMAGE_DIM / max(w, h)
        scaled = img.resize((int(w * scale), int(h * scale)), Image.LANCZOS)
        out_path = Path(tmp_dir) / (input_path.stem + "_scaled.png")
        scaled.save(str(out_path), format="PNG")
        return out_path


# The local-model Docling enrichment options this skill turns on with `--enrich`.
# NB: `enable_remote_services` and `PictureDescriptionApiOptions` (Docling's
# remote-VLM captioning path) are deliberately ABSENT — enrichment is
# local-model-only, so it can never become a covert data-egress channel inside
# Tier 2 that bypasses the Tier-3 gate (security boundary; see spec AC2).
ENRICH_OPTIONS = (
    "do_formula_enrichment",       # formulas → LaTeX
    "do_code_enrichment",          # code understanding
    "do_picture_classification",   # figure classification
    "do_picture_description",      # figure captioning (local model only)
)


def _configure_enrichment(opts: object) -> object:
    """Turn on the local-model enrichment flags on a Docling PdfPipelineOptions.

    Sets only the four ``do_*`` enrichment flags; it never touches
    ``enable_remote_services`` and never constructs a remote-VLM option class, so
    figure captioning uses Docling's *local* picture-description model. Pure
    (a function of the options object) so it is unit-testable without Docling."""
    for name in ENRICH_OPTIONS:
        setattr(opts, name, True)
    return opts


def _build_enriched_converter():
    """A Docling DocumentConverter with local-model enrichment on for PDF + image
    inputs. The remote-services path is never enabled (see ``_configure_enrichment``)."""
    from docling.datamodel.base_models import InputFormat
    from docling.datamodel.pipeline_options import PdfPipelineOptions
    from docling.document_converter import DocumentConverter, PdfFormatOption

    opts = _configure_enrichment(PdfPipelineOptions())
    fmt = PdfFormatOption(pipeline_options=opts)
    return DocumentConverter(format_options={InputFormat.PDF: fmt, InputFormat.IMAGE: fmt})


def _extract_docling(input_path: Path, *, enrich: bool = False) -> ExtractResult:
    try:
        from docling.document_converter import DocumentConverter
    except ImportError:
        raise RuntimeError(
            "docling is not installed and no Tier-0 extractor handles "
            f"'{input_path.suffix}'. Install docling (pip install docling) or "
            "convert to a Tier-0 format (PDF/Office/HTML/EPUB/CSV/ODF/EML)."
        )

    is_image = input_path.suffix.lower() in IMAGE_EXTS
    tmp_dir = None
    convert_path = input_path
    if is_image:
        tmp_dir = tempfile.mkdtemp()
        convert_path = _prescale_image(input_path, tmp_dir)
        if convert_path != input_path:
            print(f"WARNING: image pre-scaled to fit {MAX_IMAGE_DIM}px before OCR")
    try:
        # Enrichment adds local models to the pipeline; the default path stays the
        # bare converter so its Tier-2 body is byte-identical to slice 2.
        converter = _build_enriched_converter() if enrich else DocumentConverter()
        result = converter.convert(str(convert_path))
        markdown = result.document.export_to_markdown()
    finally:
        if tmp_dir:
            import shutil

            shutil.rmtree(tmp_dir, ignore_errors=True)

    # The Docling body is passed through unmodified — the builder wraps, never
    # rewrites. Enriched captions/formulas/code land here as inert body content
    # (the leading-block-only contract makes them un-forgeable), never as
    # instructions: they are model output derived from an untrusted document image.
    return ExtractResult(
        body=markdown,
        tier=contract.TIER_2,
        content_type="image" if is_image else input_path.suffix.lower().lstrip("."),
        confidence="high",
        requires_review=False,
    )


# --- Assemble + confined write ----------------------------------------------


def assemble(result: ExtractResult, source_name: str) -> str:
    """Wrap an extraction in the unified output contract (frontmatter + body).

    The frontmatter is the leading ``---``-fenced block only; the body sits
    below it, so a ``---`` line in the body is content, not a second block."""
    fields: dict[str, object] = {
        "source-file": source_name,
        "content-type": result.content_type,
        "ingestion-date": contract.now_iso(),
    }
    if result.escalation is not None:
        fields["escalation-target"] = result.escalation
    frontmatter = contract.build_frontmatter(
        tier=result.tier,
        extraction_confidence=result.confidence,
        requires_review=result.requires_review,
        fields=fields,
    )
    return frontmatter + "\n\n" + result.body.rstrip("\n") + "\n"


def write_output(input_path: Path, text: str) -> Path:
    """Write ``<basename>.md`` next to the input, confined to the input's
    resolved directory (realpath + component containment)."""
    root = input_path.resolve().parent
    output_path = safe_io.confine(root / (input_path.stem + ".md"), root)
    output_path.write_text(text, encoding="utf-8")
    return output_path


def convert_file(input_path: Path, *, enrich: bool = False) -> None:
    if not input_path.exists():
        print(f"ERROR: File not found: {input_path}", file=sys.stderr)
        sys.exit(1)
    if input_path.suffix.lower() not in supported_exts():
        print(
            f"ERROR: Unsupported file type '{input_path.suffix}'. Supported: "
            f"{', '.join(sorted(supported_exts()))}",
            file=sys.stderr,
        )
        sys.exit(1)

    try:
        result = dispatch(input_path, enrich=enrich)
    except Exception as exc:
        print(
            f"ERROR: could not convert {input_path.name}: {exc}\n"
            "If the file is password-protected or encrypted, remove the "
            "protection and retry. If it may be corrupt, confirm it opens in "
            "its native application.",
            file=sys.stderr,
        )
        sys.exit(1)

    text = assemble(result, input_path.name)
    output_path = write_output(input_path, text)

    print(f"OUTPUT: {output_path}")
    print(f"LINES: {len(text.splitlines())}")
    print(f"WORDS: {len(result.body.split())}")
    if result.requires_review:
        target = f" — escalate to Tier {result.escalation}" if result.escalation else ""
        print(f"WARNING: extraction flagged requires-review (low confidence){target}")


# --- CLI --------------------------------------------------------------------


def cmd_check(names: list[str]) -> int:
    """Probe optional Tier-0 libraries; exit 0 if all present, 2 if any absent
    (mirrors the sibling markdown-to-* --check probe)."""
    targets = names or list(OPTIONAL_LIBS)
    missing = False
    for name in targets:
        present = _lib_available(name)
        print(f"{name}: {'present' if present else 'absent'}")
        if not present:
            missing = True
            if name in OPTIONAL_LIBS:
                print(f"  install: {OPTIONAL_LIBS[name]}", file=sys.stderr)
    return 2 if missing else 0


# --- Extractor registry (defined after the extractors, before the CLI) ------
_EXTRACTORS.update({
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".xlsx": _extract_xlsx,
    ".pptx": _extract_pptx,
    ".html": _extract_html,
    ".htm": _extract_html,
    ".epub": _extract_epub,
    ".csv": _extract_csv,
    ".tsv": _extract_csv,
    ".odt": _extract_odf,
    ".ods": _extract_odf,
    ".odp": _extract_odf,
    ".eml": _extract_eml,
})


def _build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(
        prog="convert.py",
        description="Convert documents and images to Markdown, tiered and no-ML-first.",
    )
    p.add_argument("files", nargs="*",
                   help="input file(s); each is converted to <basename>.md")
    # `--check` keeps its documented positional library-name form: absent → None;
    # `--check` alone → [] (probe all optional libs); `--check pypdf` → ["pypdf"].
    # nargs="*" keeps the library names off the `files` positional.
    p.add_argument("--check", nargs="*", default=None, metavar="LIB",
                   help="probe optional Tier-0 libraries (all if none named) and exit")
    p.add_argument("--enrich", action="store_true",
                   help="Tier-2 only: turn on local-model Docling enrichment "
                        "(formulas→LaTeX, code, figure classification + captioning)")
    return p


def main(argv: list[str] | None = None) -> int:
    args = sys.argv[1:] if argv is None else argv
    ns = _build_parser().parse_args(args)

    if ns.check is not None:
        return cmd_check(ns.check)
    if not ns.files:
        print("Usage: convert.py <file> [file2 ...] [--enrich] | --check [library ...]",
              file=sys.stderr)
        return 1
    for arg in ns.files:
        convert_file(Path(arg), enrich=ns.enrich)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
