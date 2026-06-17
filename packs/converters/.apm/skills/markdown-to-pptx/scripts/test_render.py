"""TDD suite for the markdown-to-pptx renderer.

Covers the four contracts the spec names: placeholder-manifest enumeration,
Markdown→content-model mapping, the Tier-1 `--check` probe, and output-path
confinement — plus an end-to-end render that re-opens the produced .pptx and
asserts the filled values are present (the work-loop "exercise the real
artifact" requirement).

The fixture template is built from the python-pptx default presentation at test
time rather than committed as an opaque binary: the builder below is its
authorship, so a reviewer regenerates it by reading this file. Run with
`python -m pytest` from this directory.
"""

from __future__ import annotations

import subprocess
import sys
from pathlib import Path

import pytest

import render

HERE = Path(__file__).resolve().parent

MARKDOWN = """\
---
title: Q3 Review
subtitle: Prepared for the board
---

# Highlights

- Revenue up 12%
- Two new markets

# Numbers

| Metric | Value |
| --- | --- |
| ARR | 4.2M |
| Churn | 1.1% |
"""


def _build_template(path: Path) -> Path:
    """Author a minimal .pptx template (the python-pptx default master)."""
    from pptx import Presentation

    Presentation().save(str(path))
    return path


# --------------------------------------------------------------------------- #
# inspect — placeholder manifest
# --------------------------------------------------------------------------- #
def test_inspect_returns_manifest_keyed_by_idx(tmp_path):
    template = _build_template(tmp_path / "template.pptx")
    manifest = render.inspect_template(template)

    assert manifest, "default template should expose layout placeholders"
    assert all({"layout", "idx", "type", "name"} <= rec.keys() for rec in manifest)
    # The title layout exposes a title placeholder at its stable idx 0.
    title = [r for r in manifest if r["layout"] == 0 and r["idx"] == 0]
    assert title and title[0]["type"] in {"TITLE", "CENTER_TITLE"}
    # Records are keyed by placeholder idx, not list position: idx 11 (FOOTER)
    # is present even though it is not the second placeholder in the layout.
    assert any(r["idx"] == 11 for r in manifest)


# --------------------------------------------------------------------------- #
# map — Markdown → content model
# --------------------------------------------------------------------------- #
def test_empty_markdown_parses_without_crash():
    assert render.parse_markdown("") == {"scalars": {}, "sections": []}


def test_render_warns_when_layout_has_no_body_placeholder(tmp_path, monkeypatch):
    template = _build_template(tmp_path / "template.pptx")
    model = {"scalars": {}, "sections": [{"heading": "H", "bullets": ["a"], "table": None}]}
    # Simulate a layout that exposes no body/object placeholder.
    monkeypatch.setattr(render, "_placeholder_by_type", lambda *a, **k: None)
    filled, warnings = render.render_pptx(model, template, tmp_path / "out.pptx")
    assert any("body placeholder" in w and "dropped" in w for w in warnings)


def test_map_projects_frontmatter_headings_lists_and_table():
    model = render.parse_markdown(MARKDOWN)

    # front-matter → scalars (the title-slide source)
    assert model["scalars"]["title"] == "Q3 Review"
    assert model["scalars"]["subtitle"] == "Prepared for the board"

    # one H1/H2 heading → one section (→ one slide)
    headings = [s["heading"] for s in model["sections"]]
    assert headings == ["Highlights", "Numbers"]

    # a list → bullet rows
    assert model["sections"][0]["bullets"] == ["Revenue up 12%", "Two new markets"]
    assert model["sections"][0]["table"] is None

    # a Markdown table → table rows (placed into a TABLE-type placeholder at render)
    assert model["sections"][1]["table"] == [["Metric", "Value"], ["ARR", "4.2M"], ["Churn", "1.1%"]]


# --------------------------------------------------------------------------- #
# --check — Tier-1 probe
# --------------------------------------------------------------------------- #
def test_check_exit_zero_when_present():
    assert render.cmd_check() == 0


def test_check_exit_two_when_absent(monkeypatch):
    monkeypatch.setitem(sys.modules, "pptx", None)  # forces ImportError on `import pptx`
    assert render.cmd_check() == 2


# --------------------------------------------------------------------------- #
# confinement — output/template path stays under the working directory
# --------------------------------------------------------------------------- #
def test_confine_rejects_parent_traversal(tmp_path):
    root = tmp_path / "workdir"
    root.mkdir()
    with pytest.raises(ValueError):
        render.confine(root / ".." / "evil.pptx", root)


def test_confine_rejects_symlink_escape(tmp_path):
    root = tmp_path / "workdir"
    root.mkdir()
    outside = tmp_path / "outside"
    outside.mkdir()
    link = root / "link"
    link.symlink_to(outside, target_is_directory=True)
    with pytest.raises(ValueError):
        render.confine(link / "evil.pptx", root)


def test_confine_rejects_sibling_prefix(tmp_path):
    root = tmp_path / "workdir"
    root.mkdir()
    (tmp_path / "workdir-evil").mkdir()
    with pytest.raises(ValueError):
        render.confine(tmp_path / "workdir-evil" / "out.pptx", root)


def test_confine_accepts_path_under_root(tmp_path):
    root = tmp_path / "workdir"
    root.mkdir()
    assert render.confine(root / "deck.pptx", root) == (root / "deck.pptx").resolve()


# --------------------------------------------------------------------------- #
# end-to-end render (manual-QA mode, exercised as an integration test)
# --------------------------------------------------------------------------- #
def test_render_fills_template_and_reopened_values_present(tmp_path):
    from pptx import Presentation

    template = _build_template(tmp_path / "template.pptx")
    (tmp_path / "deck.md").write_text(MARKDOWN, encoding="utf-8")

    result = subprocess.run(
        [sys.executable, str(HERE / "render.py"), "render", "deck.md",
         "--template", "template.pptx", "--output", "out.pptx"],
        cwd=tmp_path, capture_output=True, text=True,
    )
    assert result.returncode == 0, result.stderr
    assert "OUTPUT:" in result.stdout
    out = tmp_path / "out.pptx"
    assert out.is_file()

    # Re-open the produced file and assert the mapped values survived the render.
    prs = Presentation(str(out))
    texts = [
        shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    ]
    blob = "\n".join(texts)
    assert "Q3 Review" in blob
    assert "Revenue up 12%" in blob

    table_cells = [
        cell.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_table
        for row in shape.table.rows
        for cell in row.cells
    ]
    assert "ARR" in table_cells and "4.2M" in table_cells
